import logging
from pathlib import Path
from pptx import Presentation
from pypdf import PdfReader
from dataclasses import dataclass

logger = logging.getLogger("study_companion")


@dataclass
class TextChunk:
    text: str
    source: str
    page: int


def extract_text_from_pptx(file_path: str, display_name: str = None) -> list[TextChunk]:
    chunks = []
    prs = Presentation(file_path)
    source = display_name or Path(file_path).name      # use display_name

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            combined = " ".join(slide_text)
            chunks.append(TextChunk(
                text=combined,
                source=source,             # now uses real filename
                page=slide_num
            ))
    return chunks


def extract_text_from_pdf(file_path: str, display_name: str = None) -> list[TextChunk]:
    chunks = []
    reader = PdfReader(file_path)
    source = display_name or Path(file_path).name      # use display_name

    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text and text.strip():
            chunks.append(TextChunk(
                text=text.strip(),
                source=source,             # now uses real filename
                page=page_num
            ))
    return chunks


def split_into_chunks(
    raw_chunks: list[TextChunk],
    chunk_size: int = 500,
    overlap: int = 50
) -> list[TextChunk]:
    final_chunks = []

    for raw in raw_chunks:
        text = raw.text
        words = text.split()

        if len(words) <= chunk_size:
            final_chunks.append(TextChunk(
                text=text,
                source=raw.source,  # Keeps it bound to your updated reference
                page=raw.page
            ))
            continue

        start = 0
        while start < len(words):
            end = start + chunk_size
            chunk_words = words[start:end]
            chunk_text = " ".join(chunk_words)

            final_chunks.append(TextChunk(
                text=chunk_text,
                source=raw.source,
                page=raw.page
            ))

            start += chunk_size - overlap

    logger.info(f"Split into {len(final_chunks)} final chunks")
    return final_chunks


def process_file(
    file_path: str,
    chunk_size: int = 500,
    overlap: int = 50,
    original_filename: str = None
) -> list[TextChunk]:
    path = Path(file_path)
    suffix = path.suffix.lower()
    display_name = original_filename or path.name 

    if suffix == ".pptx":
        raw_chunks = extract_text_from_pptx(file_path, display_name)
    elif suffix == ".pdf":
        raw_chunks = extract_text_from_pdf(file_path, display_name)
    else:
        raise ValueError(f"Unsupported file type: {suffix}. Only .pptx and .pdf supported.")

    return split_into_chunks(raw_chunks, chunk_size, overlap)